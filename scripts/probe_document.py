#!/usr/bin/env python3
"""Classify one document and pick the parsing tier for it.

Runs on CPU in well under 100 ms and loads no model. Every threshold here comes
from measurements on the BLVERA corpus; see SKILL.md for the evidence table.

Usage: probe_document.py <file> [--json]
"""
import argparse, json, pathlib, platform, sys, zipfile

# --- routing thresholds (measured, not guessed) -----------------------------
MIN_CHARS_PER_PAGE = 50      # below this a PDF has no usable text layer
TABLE_GRID_PATHS = 400       # vector paths/page above this = dense ruled table
TABLE_GRID_WARN = 100        # between WARN and GRID = some ruling, still cheap
SPARSE_TEXT_PAGE = 1200      # few chars + an image = datasheet/brochure layout
MERGED_CELLS_HEAVY = 50      # merged regions in a workbook
MERGED_TABLES_HEAVY = 3      # DOCX tables carrying gridSpan/vMerge

LEGACY_OFFICE = {".doc", ".xls", ".ppt", ".xlsb", ".rtf", ".odt", ".ods", ".odp"}
MARKITDOWN_ONLY = {".msg", ".ipynb", ".zip", ".epub", ".mp3", ".wav", ".m4a", ".csv", ".json", ".xml"}


def probe_pdf(path, sample=8):
    import pypdfium2 as pdfium
    doc = pdfium.PdfDocument(path)
    n = len(doc)
    s = min(n, sample)
    chars = paths = images = 0
    for i in range(s):
        page = doc[i]
        chars += len(page.get_textpage().get_text_range())
        for o in page.get_objects():
            if o.type == 2:
                paths += 1
            elif o.type == 3:
                images += 1
    return dict(pages=n, chars_per_page=round(chars / s, 1),
                paths_per_page=round(paths / s, 1), images_per_page=round(images / s, 2))


def probe_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True, read_only=False)
    hidden = [ws.title for ws in wb.worksheets if ws.sheet_state != "visible"]
    merged = sum(len(ws.merged_cells.ranges) for ws in wb.worksheets)
    return dict(sheets=len(wb.worksheets), hidden_sheets=len(hidden),
                hidden_sheet_names=hidden, merged_cells=merged)


def probe_docx(path):
    import docx
    d = docx.Document(path)
    headings = sum(1 for p in d.paragraphs
                   if p.style and p.style.name.startswith("Heading"))
    merged_tables = sum(1 for t in d.tables
                        if "gridSpan" in t._tbl.xml or "vMerge" in t._tbl.xml)
    raw = zipfile.ZipFile(path).read("word/document.xml").decode("utf-8", "ignore")
    return dict(paragraphs=len(d.paragraphs), heading_styles=headings,
                tables=len(d.tables), tables_with_merge=merged_tables,
                gridspan=raw.count("<w:gridSpan"), vmerge=raw.count("<w:vMerge"))


def probe_pptx(path):
    from pptx import Presentation
    p = Presentation(path)
    titled = sum(1 for s in p.slides if s.shapes.title is not None
                 and s.shapes.title.text.strip())
    return dict(slides=len(p.slides), slides_with_title_placeholder=titled)


def default_ocr_engine():
    """macOS Vision measured best on Vietnamese; Tesseract is the portable pick."""
    return ("ocrmac", "vi-VT") if platform.system() == "Darwin" else ("tesseract", "vie")


def route(info):
    """Return (tier, engine, reason, flags) from the probe facts."""
    ext, k = info["ext"], info["kind"]
    flags = []

    if ext in MARKITDOWN_ONLY:
        return "T0b", "markitdown", "định dạng chỉ markitdown đọc được", flags
    if ext in LEGACY_OFFICE:
        return "T0", "anydoc", "định dạng Office nhị phân cũ, chỉ anydoc đọc", flags

    if k == "xlsx":
        if info["hidden_sheets"] > 0:
            flags.append("HIDDEN_SHEET_LEAK_RISK")
            return "T1", "mineru", f"{info['hidden_sheets']} sheet ẩn — engine rẻ sẽ lộ ra", flags
        if info["merged_cells"] > MERGED_CELLS_HEAVY:
            return "T1", "mineru", f"{info['merged_cells']} vùng gộp ô — cần HTML giữ span", flags
        return "T0", "anydoc", "workbook đơn giản", flags

    if k == "docx":
        if info["heading_styles"] == 0 and info["paragraphs"] > 20:
            flags.append("NO_HEADING_STYLES")
        if info["tables_with_merge"] >= MERGED_TABLES_HEAVY:
            return "T1", "mineru", f"{info['tables_with_merge']} bảng có ô gộp", flags
        return "T0", "anydoc", "văn bản phẳng, bảng đơn giản", flags

    if k == "pptx":
        return "T1", "mineru", "chỉ mineru dựng được tiêu đề slide", flags

    if k == "pdf":
        eng, lang = default_ocr_engine()
        if info["chars_per_page"] < MIN_CHARS_PER_PAGE:
            flags.append("NEEDS_OCR")
            return "T2", f"docling:ocr:{eng}:{lang}", "PDF không có text layer", flags
        if info["paths_per_page"] > TABLE_GRID_PATHS:
            flags.append("DENSE_TABLE_GRID")
            return "T2", "docling", f"{info['paths_per_page']:.0f} nét vẽ/trang — bảng lưới dày", flags
        if info["images_per_page"] >= 0.5 and info["chars_per_page"] < SPARSE_TEXT_PAGE:
            flags.append("BORDERLESS_SPEC_TABLE")
            return "T2", "docling", "ít chữ + có hình — nghi bảng thông số không kẻ viền", flags
        if info["paths_per_page"] > TABLE_GRID_WARN:
            flags.append("LAYOUT_RISK_MEDIUM")
        return "T0", "anydoc", "PDF text layer, bố cục đơn giản", flags

    return "T0b", "markitdown", "loại không nhận diện được — thử engine phủ rộng nhất", flags


def probe(path):
    path = pathlib.Path(path)
    ext = path.suffix.lower()
    kind = {".pdf": "pdf", ".xlsx": "xlsx", ".xlsm": "xlsx",
            ".docx": "docx", ".pptx": "pptx"}.get(ext, "other")
    info = dict(file=str(path), ext=ext, kind=kind, size_bytes=path.stat().st_size)
    try:
        if kind == "pdf":
            info.update(probe_pdf(path))
        elif kind == "xlsx":
            info.update(probe_xlsx(path))
        elif kind == "docx":
            info.update(probe_docx(path))
        elif kind == "pptx":
            info.update(probe_pptx(path))
    except Exception as e:
        info["probe_error"] = f"{type(e).__name__}: {e}"
    tier, engine, reason, flags = route(info)
    info.update(tier=tier, engine=engine, reason=reason, flags=flags)
    if kind == "pdf" and info.get("chars_per_page", 1) < MIN_CHARS_PER_PAGE:
        info["kind"] = "pdf_scan"
    return info


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("file")
    ap.add_argument("--json", action="store_true")
    a = ap.parse_args()
    info = probe(a.file)
    if a.json:
        print(json.dumps(info, ensure_ascii=False, indent=1))
    else:
        print(f"{pathlib.Path(a.file).name}")
        print(f"  loại      : {info['kind']}")
        for k in ("pages", "chars_per_page", "paths_per_page", "images_per_page",
                  "sheets", "hidden_sheets", "merged_cells", "heading_styles",
                  "tables", "tables_with_merge", "slides", "slides_with_title_placeholder"):
            if k in info:
                print(f"  {k:10s}: {info[k]}")
        print(f"  → {info['tier']} / {info['engine']}  ({info['reason']})")
        if info["flags"]:
            print(f"  cờ        : {', '.join(info['flags'])}")


if __name__ == "__main__":
    main()
